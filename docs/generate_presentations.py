"""Generate two PowerPoint presentations for PLC4X Manager."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

RED = RGBColor(0xC8, 0x10, 0x2E)
DARK = RGBColor(0x1B, 0x1F, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
BLUE = RGBColor(0x38, 0xBD, 0xF8)

def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_red_bar(slide, prs):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf

def add_table(slide, rows_data, headers, left=0.8, top=1.8, width=11.5):
    rows = len(rows_data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(4)).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED if i > 0 else DARK
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
    return tbl


# ============================================================
# PRESENTATION 1: EXECUTIVE
# ============================================================
def generate_executive():
    prs = new_prs()

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = DARK
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
    add_text(s, 0.5, 2.7, 12, 1, "PLC4X Manager", 48, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, 0.5, 3.6, 12, 1, "Industrial Monitoring Platform", 28, False, WHITE, PP_ALIGN.CENTER)
    add_text(s, 0.5, 4.3, 12, 0.5, "Intelligent Monitoring  |  Zero License Cost  |  AI-Powered  |  Cloud-Ready", 16, False, RGBColor(0xFF,0xCC,0xCC), PP_ALIGN.CENTER)
    add_text(s, 0.5, 6.5, 12, 0.5, "Executive Presentation | 2026", 14, False, GRAY, PP_ALIGN.CENTER)

    # The Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "The Challenge", 36, True, RED)
    add_text(s, 0.8, 1.3, 12, 0.5, "Current monitoring systems are expensive, fragmented, and lack intelligence", 18, False, GRAY)
    add_bullets(s, 0.8, 2.2, 11, 4, [
        "X  Commercial SCADA licenses: $30,000-$100,000+/year per site",
        "X  Separate tools for monitoring, alarms, reports, dashboards",
        "X  No AI/ML predictions - reactive instead of proactive",
        "X  Vendor lock-in - proprietary protocols, expensive upgrades",
        "X  Long deployment time - weeks of installation and configuration",
        "X  No centralized multi-plant visibility from one screen",
    ], 20)

    # The Solution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "The Solution: PLC4X Manager", 36, True, RED)
    add_text(s, 0.8, 1.3, 12, 0.5, "One platform. All plants. Zero license cost.", 20, True, DARK)
    add_bullets(s, 0.8, 2.0, 11, 5, [
        ">> Open-source - Apache 2.0 license, no recurring fees",
        ">> Multi-protocol - Siemens S7, Modbus, OPC-UA, Allen-Bradley, 10+ protocols",
        ">> AI-powered - 7 ML algorithms: Prophet, Isolation Forest, PyOD, SHAP, and more",
        ">> One-command deploy - Docker: 'docker compose up' and it runs",
        ">> Cloud-ready - Coolify, Docker Swarm, or bare metal deployment",
        ">> Multi-plant dashboard - all sites on one screen",
        ">> HMI / Synoptic - drag-and-drop visual editor with live data & PLC write controls",
        ">> Alarms with sound - conditional thresholds per product, virtual tags",
        ">> OEE dashboard - real-time Availability x Performance x Quality",
        ">> Rockwell native - pylogix tag discovery, diagnostics, and direct reads",
    ], 17)

    # Cost Comparison
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Cost Comparison (5-Year TCO)", 36, True, RED)
    add_table(s, [
        ["License (5 years)", "$150,000+", "$75,000-$250,000", "$0"],
        ["Implementation", "$30,000-$80,000", "$50,000-$150,000", "$5,000-$15,000"],
        ["Annual maintenance", "$15,000/year", "$20,000/year", "$0"],
        ["AI/ML module", "Not available", "$30,000+", "Included (7 algorithms)"],
        ["Cloud deployment", "Extra license", "Extra license", "Included (Coolify)"],
        ["5-Year Total", "$255,000+", "$225,000-$550,000", "$5,000-$15,000"],
    ], ["", "Ignition (Inductive)", "WinCC / AVEVA", "PLC4X Manager"])
    add_text(s, 0.8, 6.5, 11, 0.5, "Savings: $200,000-$500,000+ over 5 years per site", 20, True, GREEN, PP_ALIGN.CENTER)

    # AI & Intelligence
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Artificial Intelligence Built-In", 36, True, RED)
    add_text(s, 0.8, 1.3, 12, 0.5, "From reactive monitoring to predictive intelligence — 7 ML algorithms", 18, False, GRAY)
    add_bullets(s, 0.8, 2.0, 5.5, 4, [
        "Prophet Time-Series Forecasting",
        "  - Predicts tag values 2-24h ahead",
        "  - Confidence intervals (upper/lower)",
        "",
        "Isolation Forest + PyOD Ensemble",
        "  - 3-algorithm consensus anomaly",
        "  - ECOD + LOF + Isolation Forest",
        "",
        "SHAP Explainability",
        "  - Top contributing tags per anomaly",
        "  - Human-readable explanations",
    ], 16)
    add_bullets(s, 6.8, 2.0, 5.5, 4, [
        "Change Point Detection (ruptures)",
        "  - Detects sudden mean/trend shifts",
        "  - Real-time process drift alerts",
        "",
        "Pattern Matching (stumpy)",
        "  - Recurring patterns (motifs)",
        "  - Rare anomalies (discords)",
        "",
        "Cross-Tag Correlation",
        "  - Detects broken correlations",
        "  - Process relationship monitoring",
    ], 16)
    add_text(s, 0.8, 6.0, 11, 1, "Result: Predict equipment failures before they happen.\nReduce unplanned downtime by up to 30%.", 20, True, DARK, PP_ALIGN.CENTER)

    # Rockwell Integration
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Rockwell Allen-Bradley Native Integration", 36, True, RED)
    add_text(s, 0.8, 1.3, 12, 0.5, "PLC4X Manager monitors — Rockwell operates. Direct EtherNet/IP via pylogix.", 18, False, GRAY)
    add_bullets(s, 0.8, 2.2, 5.5, 4, [
        "Automatic Tag Discovery",
        "  - Controller-scoped tags",
        "  - Program-scoped tags",
        "  - Full tag inventory in seconds",
        "",
        "PLC Diagnostics",
        "  - Model identification (1768-ENBT/A)",
        "  - Firmware version detection",
        "  - Serial number & vendor ID",
        "  - Health check (tags + programs)",
    ], 16)
    add_bullets(s, 6.8, 2.2, 5.5, 4, [
        "Safety-First Design",
        "  - Global read-only mode (default ON)",
        "  - One CIP connection at a time",
        "  - 30s rate limit between connections",
        "  - 10s max operation timeout",
        "  - No retry on failure (fail-fast)",
        "",
        "Supported Types",
        "  - BOOL, SINT, INT, DINT, LINT, REAL",
        "  - STRING, arrays, UDTs",
        "  - CompactLogix & ControlLogix",
    ], 16)

    # ROI
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Return on Investment", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 11, 4, [
        "License savings: $30,000-$100,000/year per site",
        "Reduced downtime: 10-30% fewer unplanned stops (AI predictions)",
        "Energy optimization: 5-15% savings from real-time monitoring",
        "Compliance: full audit trail for regulatory requirements (SQLite DB)",
        "Faster troubleshooting: instant tag trending, alarm history",
        "Reduced tool sprawl: one platform instead of multiple tools",
        "Fast deployment: running in hours, not weeks (Docker or Coolify)",
        "20-year lifespan: open-source, no vendor sunset risk",
    ], 20)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5), Inches(9), Inches(1.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Payback period: < 3 months"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    # Features
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Complete Feature Set", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Multi-plant dashboard",
        "HMI / Synoptic editor (16 elements)",
        "Real-time alarm system with sound",
        "OEE monitoring (A x P x Q)",
        "PDF reports & CSV export",
        "Shift logbook",
        "Inline tag trending",
        "Rockwell tag discovery & diagnostics",
        "WebSocket real-time data push",
    ], 17)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "7 ML algorithms (Prophet, PyOD, SHAP...)",
        "8 Grafana dashboards pre-configured",
        "RBAC (admin / operator / monitor)",
        "Plant-based access filter per user",
        "Full audit trail (SQLite DB)",
        "MQTT real-time streaming",
        "Kiosk mode for TV displays",
        "Virtual & calculated tags",
        "Coolify / cloud deployment ready",
    ], 17)
    add_text(s, 0.8, 6.0, 11, 0.5, "10+ industrial protocols  |  6 Docker containers  |  Open source (Apache 2.0)", 16, False, GRAY, PP_ALIGN.CENTER)

    # Next Steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Next Steps", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 11, 4, [
        "1. Pilot: deploy at one plant with existing PLCs (1-2 weeks)",
        "2. Validate: run in parallel with current system for 30 days",
        "3. Expand: roll out to all plants (1 day per plant)",
        "4. Optimize: enable AI predictions, configure OEE targets",
        "5. Scale: add Grafana dashboards for management visibility",
        "6. Cloud: deploy to Coolify for centralized multi-site management",
    ], 22)
    add_text(s, 0.8, 5.5, 11, 1, "No risk: runs alongside existing systems.\nNo disruption to production.", 20, True, DARK, PP_ALIGN.CENTER)

    import os; prs.save(os.path.join(os.path.dirname(__file__), "PLC4X_Manager_Executive_Presentation.pptx"))
    print("Executive presentation: DONE")


# ============================================================
# PRESENTATION 2: IT/TA TECHNICAL
# ============================================================
def generate_technical():
    prs = new_prs()

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = DARK
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
    add_text(s, 0.5, 2.7, 12, 1, "PLC4X Manager", 48, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, 0.5, 3.6, 12, 1, "Technical Architecture & Implementation Guide", 24, False, WHITE, PP_ALIGN.CENTER)
    add_text(s, 0.5, 4.3, 12, 0.5, "For IT Infrastructure & Industrial Automation Teams", 16, False, RGBColor(0xFF,0xCC,0xCC), PP_ALIGN.CENTER)
    add_text(s, 0.5, 6.5, 12, 0.5, "Technical Presentation | 2026", 14, False, GRAY, PP_ALIGN.CENTER)

    # Architecture Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "System Architecture", 36, True, RED)
    add_text(s, 0.8, 1.2, 12, 0.5, "6 Docker containers, fully orchestrated with docker-compose", 18, False, GRAY)
    add_table(s, [
        ["plc4x-admin", "3080, 3443", "FastAPI + uvicorn, Web UI, HMI editor, WebSocket, MQTT/InfluxDB writer, SQLite DB"],
        ["plc4x-server", "12687", "Apache PLC4X OPC-UA Server (Java) + pylogix service (Python, port 5000)"],
        ["influxdb", "8086", "Time-series DB: 90-day raw, 2-year hourly, daily forever"],
        ["grafana", "3000", "8 dashboards, 5 alert rules, InfluxDB + MQTT datasources"],
        ["mosquitto", "1883", "MQTT broker for real-time tag streaming"],
        ["plc4x-ml", "-", "7 ML algorithms: Prophet, IForest, PyOD, SHAP, ruptures, stumpy, correlation"],
    ], ["Container", "Ports", "Description"], top=1.8)

    # Data Flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Data Flow Pipeline", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 11, 5, [
        "1. PLC4X Server connects to PLCs via S7/Modbus/OPC-UA/EIP protocols",
        "2. PLC4X Server exposes all tags via OPC-UA (port 12687)",
        "3. pylogix service (port 5000) provides direct Rockwell tag discovery & diagnostics",
        "4. Background Poller (Python asyncua) reads OPC-UA every 5s (configurable per device)",
        "5. Poller writes to: JSON cache (instant API), InfluxDB (SYNCHRONOUS), MQTT (real-time)",
        "6. Virtual tags bypass OPC-UA — stored in JSON file, written via REST API",
        "7. Calculated tags evaluated after each poll: formulas on other tag values",
        "8. Alarm engine evaluates thresholds every cycle, persists to SQLite alarm_history table",
        "9. WebSocket pushes live data, alarm events, and MQTT messages to browser clients",
        "10. ML container queries InfluxDB every 5min: 7 algorithms run on all active tags",
        "11. Grafana queries InfluxDB directly for dashboards and alerts",
        "12. Web UI reads from JSON cache (~60ms) instead of direct OPC-UA (~1s)",
    ], 15)

    # Supported Protocols
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Supported PLC Protocols", 36, True, RED)
    add_table(s, [
        ["Siemens S7", "s7://192.168.1.10", "%DB1:0:REAL"],
        ["Modbus TCP", "modbus-tcp://192.168.1.20:502", "holding-register:1"],
        ["Modbus RTU/ASCII", "modbus-rtu:///dev/ttyUSB0", "holding-register:1"],
        ["OPC-UA Client", "opcua:tcp://192.168.1.30:4840", "ns=2;i=10"],
        ["Allen-Bradley Logix", "logix://192.168.1.50", "MyTag"],
        ["EtherNet/IP", "eip://192.168.1.40", "%MyTag:DINT"],
        ["KNXnet/IP", "knxnet-ip://192.168.1.60", "1/1/1:DPT_Switch"],
        ["IEC 60870-5-104", "iec-60870-5-104://192.168.1.70", "M_ME_NC_1:5:20"],
        ["Firmata (Arduino)", "firmata:///dev/ttyACM0", "digital:13"],
        ["Simulated (Test)", "simulated://127.0.0.1", "RANDOM/Temporary:DINT"],
    ], ["Protocol", "Connection String", "Tag Example"])

    # Rockwell / pylogix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Rockwell Allen-Bradley Integration (pylogix)", 36, True, RED)
    add_text(s, 0.8, 1.2, 12, 0.5, "Direct EtherNet/IP CIP communication — no RSLinx or middleware required", 18, False, GRAY)
    add_bullets(s, 0.8, 1.8, 5.5, 5, [
        "Discovery Endpoints:",
        "  POST /api/plctag/discover",
        "    - Controller + program tags",
        "    - Tag name, type, array, UDT info",
        "  POST /api/plctag/discover/programs",
        "    - List all program names",
        "",
        "Diagnostics Endpoints:",
        "  POST /api/plctag/diagnostics/identity",
        "    - Vendor, model, firmware, serial",
        "  POST /api/plctag/diagnostics/health",
        "    - Identity + tag count + programs",
    ], 14)
    add_bullets(s, 6.8, 1.8, 5.5, 5, [
        "Safety Rules (plctag_safety.py):",
        "  R1: Global lock - ONE CIP connection",
        "  R2: Short-lived - 10s max per operation",
        "  R3: Socket timeout - 3-5s (no hung TCP)",
        "  R5: Rate limit - 30s cooldown per PLC IP",
        "  R7: Connection counting & logging",
        "  R8: No retry - fail-fast on error",
        "  R10: Admin-only access",
        "",
        "Supported PLC Types:",
        "  - CompactLogix (1769-L series)",
        "  - ControlLogix (1756-L series)",
        "  - Communication via 1768-ENBT/A",
    ], 14)

    # Security Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Security Architecture", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Authentication:",
        "  - JWT HS256 tokens (configurable expiry)",
        "  - Brute-force: 5 fails = 5-min IP lockout",
        "  - API key for machine-to-machine access",
        "",
        "Authorization (RBAC):",
        "  - 3 roles: admin / operator / monitor",
        "  - Plant-based access filter (USERS_JSON)",
        "  - Per-device allowWrite flag",
        "",
        "PLC Write Protection (4 layers):",
        "  - PLC_READONLY=true (global, default ON)",
        "  - Device allowWrite=false (per device)",
        "  - Role 'monitor' (no PLC writes)",
        "  - Plant filter (restrict by plant)",
    ], 14)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "Transport:",
        "  - HTTPS auto-generated TLS (port 3443)",
        "  - OPC-UA Basic256Sha256 + SignAndEncrypt",
        "  - Certificate-based auth (PKI trust store)",
        "",
        "Data Protection:",
        "  - SQLite audit trail (user, IP, timestamp)",
        "  - Flux injection regex allowlist",
        "  - XSS protection (escHtml/escAttr)",
        "  - HMAC constant-time comparison",
        "  - FileLock cross-process file access",
        "",
        "PLC Safety:",
        "  - Global CIP connection lock",
        "  - Rate limiting per PLC IP (30s)",
        "  - Socket timeout (prevents hung connections)",
    ], 14)

    # Alarm System
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Alarm System - Deep Dive", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Per-tag thresholds:",
        "  - warningHigh / warningLow",
        "  - criticalHigh / criticalLow",
        "  - Validation: CL <= WL <= WH <= CH",
        "",
        "Conditional profiles:",
        "  - Condition tag (e.g., ProductCode)",
        "  - Multiple profiles per product",
        "  - Fallback to default thresholds",
        "",
        "Sound notification:",
        "  - 880Hz beep for critical unack",
        "  - 30s cooldown between beeps",
        "  - AudioContext with resume()",
    ], 14)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "Evaluation engine:",
        "  - Runs every poller cycle",
        "  - Persisted to SQLite alarm_history",
        "  - Active alarms in .alarms.json",
        "",
        "State management:",
        "  - Active alarms with acknowledge",
        "  - History (SQLite, unlimited)",
        "  - Severity escalation tracking",
        "  - Orphan cleanup on config change",
        "",
        "Integration:",
        "  - WebSocket: real-time alarm push",
        "  - MQTT: plc4x/_alarms/{device}/{tag}",
        "  - Grafana: 5 pre-configured alert rules",
        "  - HMI: alarm banner element (blinks)",
    ], 14)

    # HMI System
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "HMI / Synoptic System", 36, True, RED)
    add_text(s, 0.8, 1.2, 12, 0.5, "Konva.js 9.x canvas with 16 element types", 18, False, GRAY)
    add_bullets(s, 0.8, 1.8, 5.5, 5, [
        "Hierarchy: Plant > Area > Equipment > Screen",
        "Full CRUD + rename for all levels",
        "Drag-and-drop element placement",
        "Property panel with tag selection",
        "Per-element device override",
        "Background image upload",
        "Zoom/pan with mouse wheel",
        "Fullscreen mode (?hmi=equip-id)",
        "Auto-save on every change",
        "Demo with 21 elements pre-configured",
    ], 15)
    add_bullets(s, 6.8, 1.8, 5.5, 5, [
        "16 Element Types:",
        "  Display, Gauge, Tank, Bar Graph,",
        "  Progress Bar, Indicator, Valve, Motor,",
        "  Pipe (animated flow), Alarm Banner,",
        "  Label, Image, Button, Switch,",
        "  Slider, Numeric Input",
        "",
        "Write Controls (with confirmation):",
        "  - Button: write value on click",
        "  - Switch: toggle ON/OFF",
        "  - Valve: open/close",
        "  - Motor: start/stop",
        "  - Slider: drag to set value",
        "  - Numeric Input: type value",
    ], 14)

    # Virtual & Calculated Tags
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Virtual & Calculated Tags", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Virtual Tags (address: VIRTUAL):",
        "  - No PLC connection needed",
        "  - Value set via REST API or HMI",
        "  - Stored in .virtual-tags.json",
        "  - Cross-process safe (FileLock)",
        "  - Flows to: MQTT, InfluxDB, Grafana",
        "  - Use cases:",
        "    - Product code selection",
        "    - Manual quality inputs",
        "    - Shift/batch identifiers",
        "    - External system integration",
    ], 15)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "Calculated Tags (formula-based):",
        "  - Safe expression evaluator (AST)",
        "  - Supports: +, -, *, /, **, %",
        "  - Functions: abs, min, max, sqrt, round",
        "  - Comparisons: >, <, ==, if/else",
        "  - Can chain (tag B uses tag A)",
        "  - Use cases:",
        "    - Efficiency = output / input * 100",
        "    - Delta = current - previous",
        "    - Weighted average",
        "    - Boolean logic combinations",
    ], 15)

    # OEE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "OEE (Overall Equipment Effectiveness)", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 11, 5, [
        "OEE = Availability x Performance x Quality",
        "",
        "Configuration per device:",
        "  - Running Tag: boolean/status tag indicating machine is ON",
        "  - Production Count Tag: cumulative counter (handles rollover)",
        "  - Reject Count Tag: cumulative reject counter (optional)",
        "  - Ideal Cycle Time: seconds per unit at max speed",
        "  - Planned Hours/Day: scheduled production hours",
        "",
        "Calculation from InfluxDB history:",
        "  - Availability = (running intervals sum) / planned time",
        "  - Performance = (ideal cycle time x total count) / run time",
        "  - Quality = (total count - rejects) / total count",
        "",
        "SVG gauges: green >= 85%, yellow 60-85%, red < 60%",
        "Trend chart: availability over time (hourly/daily buckets)",
        "Inline OEE on multi-plant dashboard for each configured device",
    ], 14)

    # API Reference
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "REST API Overview (50+ endpoints)", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Authentication:",
        "  POST /api/auth/login",
        "  GET /api/auth/verify",
        "",
        "Devices & Tags:",
        "  CRUD /api/devices, /api/devices/:name/tags",
        "  PUT /api/devices/:name/tags/:alias/alarms",
        "  PUT /api/devices/:name/oee-config",
        "",
        "Live Data:",
        "  GET /api/live/read",
        "  POST /api/live/write",
        "",
        "HMI:",
        "  CRUD /api/hmi/plants, areas, equipment",
        "  PUT /api/hmi/equipment/:id/screen",
    ], 13)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "Rockwell/pylogix:",
        "  POST /api/plctag/discover",
        "  POST /api/plctag/diagnostics/identity",
        "  POST /api/plctag/diagnostics/health",
        "",
        "Alarms & OEE:",
        "  GET /api/alarms",
        "  POST /api/alarms/acknowledge",
        "  GET /api/oee/calculate, /api/oee/trend",
        "",
        "Reports & Audit:",
        "  GET /api/export/csv, /api/export/pdf",
        "  GET/POST /api/logbook",
        "  GET /api/audit",
        "",
        "WebSocket: ws://.../ws/live",
    ], 13)

    # Deployment
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Deployment & Infrastructure", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "Requirements:",
        "  - Docker + Docker Compose",
        "  - 8GB RAM (16GB with ML analytics)",
        "  - Linux, Windows, or macOS",
        "  - Network access to PLCs",
        "",
        "Installation (3 commands):",
        "  git clone <repo>",
        "  cp .env.example .env",
        "  docker compose up -d",
        "",
        "Cloud Deployment (Coolify):",
        "  - GitHub auto-deploy on push",
        "  - No bind mounts (all Dockerfiles)",
        "  - Ports: 3080/3443 (no Traefik conflict)",
        "  - Named volumes for persistence",
    ], 14)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "Persistence (Docker volumes):",
        "  - config-data: YAML, SQLite, backups",
        "  - security-data: certs, keystore, PKI",
        "  - influxdb-data: time-series history",
        "  - grafana-data: dashboard customizations",
        "  - mosquitto-data: MQTT persistence",
        "",
        "Monitoring room:",
        "  - Kiosk mode: ?kiosk URL parameter",
        "  - Auto-cycles: Dashboard, HMI, OEE",
        "  - 20 kiosks, 60 clients supported",
        "  - Fullscreen browser on TV",
        "",
        "Networking:",
        "  - Bridge mode (default) or host mode",
        "  - HTTPS on port 3443 (self-signed)",
    ], 14)

    # Grafana & ML
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Grafana Dashboards & ML Analytics", 36, True, RED)
    add_bullets(s, 0.8, 1.5, 5.5, 5, [
        "8 Pre-configured Dashboards:",
        "  1. Overview - all devices at a glance",
        "  2. Plant Overview - filtered by plant",
        "  3. Device Detail - deep dive per device",
        "  4. Energy Dashboard - power metrics",
        "  5. AI Predictions - Prophet forecasts",
        "  6. Alarms & Anomalies - ML alerts",
        "  7. Custom Data - external sources",
        "  8. System Health - latency, errors",
        "",
        "5 Alert Rules:",
        "  - Device Offline (1m threshold)",
        "  - High Read Latency (>2000ms)",
        "  - ML Anomaly Detected",
        "  - No Data from Poller (2m)",
        "  - High Tag Error Rate (>10)",
    ], 14)
    add_bullets(s, 6.8, 1.5, 5.5, 5, [
        "7 ML Algorithms (plc4x-ml container):",
        "",
        "  1. Prophet - 2h forecast + confidence",
        "  2. Isolation Forest - anomaly scoring",
        "  3. PyOD Ensemble - ECOD+LOF+IForest",
        "  4. SHAP - explainable anomaly causes",
        "  5. ruptures - change point detection",
        "  6. stumpy - pattern matching (motifs)",
        "  7. numpy - cross-tag correlation",
        "",
        "Configuration (per device in config.yml):",
        "  - Enable/disable each algorithm",
        "  - Contamination, min agreement",
        "  - Segment size, penalty, window size",
        "  - Manual 'Run Now' button in UI",
    ], 14)

    # Data Retention
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_red_bar(s, prs)
    add_text(s, 0.8, 0.5, 12, 0.7, "Data Retention & Storage", 36, True, RED)
    add_table(s, [
        ["plc4x_raw", "Raw tag values", "Every poll cycle (5s default)", "90 days"],
        ["plc4x_hourly", "Hourly aggregates", "Mean, min, max per hour", "2 years"],
        ["plc4x_daily", "Daily aggregates", "Mean, min, max per day", "Forever"],
        ["plc4x_health", "Device health", "Status, latency, tag counts", "90 days"],
        ["plc4x_forecast", "ML predictions", "Prophet forecasts", "30 days"],
        ["plc4x_anomaly", "Anomaly scores", "All 7 ML algorithm results", "30 days"],
        ["SQLite: audit_entries", "Audit trail", "All write operations (user, IP)", "Unlimited"],
        ["SQLite: alarm_history", "Alarm events", "Fire, ack, clear with duration", "Unlimited"],
        ["SQLite: logbook_entries", "Shift logbook", "Operator observations", "Unlimited"],
        ["SQLite: write_log", "PLC writes", "Tag, value, user, timestamp", "Unlimited"],
    ], ["Storage", "Content", "Frequency", "Retention"])

    import os; prs.save(os.path.join(os.path.dirname(__file__), "PLC4X_Manager_Technical_Presentation.pptx"))
    print("Technical presentation: DONE")


if __name__ == "__main__":
    generate_executive()
    generate_technical()
